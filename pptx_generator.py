"""
Server-side PPTX generation (§7) — mirrors the prototype's in-browser
Present mode content, but as a real downloadable file, AAR-branded.

Uses python-pptx rather than pptxgenjs since this runs as a backend
service (no Node runtime assumed in the export worker), not an
interactively-authored one-off deck. Known python-pptx limitation kept in
mind here: setting `text_frame.text` collapses formatting, so every text
assignment below goes through a run instead.
"""
from datetime import datetime
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .brand import COLOR_ORANGE, COLOR_BLACK, COLOR_WHITE, FONT_HEADING, FONT_BODY

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _set_run(text_frame, text: str, size: int, bold: bool = False,
             color: str = COLOR_BLACK, font: str = FONT_BODY, align=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)


def _add_textbox(slide, left, top, width, height, text, size, bold=False,
                  color=COLOR_BLACK, font=FONT_BODY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    _set_run(box.text_frame, text, size, bold, color, font, align)
    return box


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def build_presentation(
    session_name: str,
    kpis: Dict[str, Any],
    top_categories: List[Dict[str, Any]],
    priority_flags: List[Dict[str, Any]],
    recommendations: List[str],
    output_path: str,
):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _title_slide(prs, session_name)
    _kpi_slide(prs, kpis)
    _top_category_slide(prs, top_categories)
    _priority_review_slide(prs, priority_flags)
    _recommendations_slide(prs, recommendations)

    prs.save(output_path)
    return output_path


def _title_slide(prs, session_name: str):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(COLOR_WHITE)
    bg.line.fill.background()

    _add_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.0),
                 "Claims Forensic Audit — AAR Insurance Kenya", 40, bold=True,
                 color=COLOR_BLACK, font=FONT_HEADING)
    _add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.6),
                 session_name, 20, color=COLOR_ORANGE, font=FONT_BODY)
    _add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
                 datetime.utcnow().strftime("Generated %d %B %Y"), 12,
                 color=COLOR_BLACK, font=FONT_BODY)


def _kpi_slide(prs, kpis: Dict[str, Any]):
    slide = _blank_slide(prs)
    _add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                 "Key Metrics", 28, bold=True, color=COLOR_BLACK, font=FONT_HEADING)

    card_w = Inches(2.9)
    card_h = Inches(1.8)
    gap = Inches(0.3)
    start_x = Inches(0.6)
    start_y = Inches(1.5)

    items = list(kpis.items())[:4]
    for i, (label, value) in enumerate(items):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb("F2F2F2")
        card.line.fill.background()

        _add_textbox(slide, x + Inches(0.15), start_y + Inches(0.2), card_w - Inches(0.3), Inches(0.8),
                     str(value), 32, bold=True, color=COLOR_ORANGE, font=FONT_HEADING)
        _add_textbox(slide, x + Inches(0.15), start_y + Inches(1.1), card_w - Inches(0.3), Inches(0.5),
                     label, 13, color=COLOR_BLACK, font=FONT_BODY)


def _top_category_slide(prs, top_categories: List[Dict[str, Any]]):
    slide = _blank_slide(prs)
    _add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                 "Top Flagged Categories", 28, bold=True, color=COLOR_BLACK, font=FONT_HEADING)

    top_y = Inches(1.6)
    row_h = Inches(0.6)
    max_amount = max((c.get("amount", 0) for c in top_categories), default=1) or 1

    for i, cat in enumerate(top_categories[:8]):
        y = top_y + i * row_h
        _add_textbox(slide, Inches(0.6), y, Inches(3.0), row_h, cat["category"], 14, color=COLOR_BLACK)
        bar_max_w = Inches(7.5)
        bar_w = Inches(max(0.1, (cat.get("amount", 0) / max_amount) * 7.5))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.8), y + Inches(0.1), bar_w, Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(COLOR_ORANGE)
        bar.line.fill.background()
        _add_textbox(slide, Inches(11.5), y, Inches(1.5), row_h,
                     f"{cat.get('flag_count', 0)}", 14, bold=True, color=COLOR_BLACK)


def _priority_review_slide(prs, priority_flags: List[Dict[str, Any]]):
    slide = _blank_slide(prs)
    _add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                 "Priority Review Items", 28, bold=True, color=COLOR_BLACK, font=FONT_HEADING)

    rows = min(len(priority_flags), 6) + 1
    table_shape = slide.shapes.add_table(rows, 3, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.0))
    table = table_shape.table
    for c, header in enumerate(["Flag Type", "Detail", "Amount / Score"]):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(COLOR_WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(COLOR_BLACK)

    for i, flag in enumerate(priority_flags[:6], start=1):
        table.cell(i, 0).text = flag.get("flag_type", "")
        table.cell(i, 1).text = str(flag.get("detail_summary", ""))
        table.cell(i, 2).text = str(flag.get("amount_or_score", ""))


def _recommendations_slide(prs, recommendations: List[str]):
    slide = _blank_slide(prs)
    _add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
                 "Recommendations", 28, bold=True, color=COLOR_BLACK, font=FONT_HEADING)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, rec in enumerate(recommendations):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {rec}"
        run.font.size = Pt(18)
        run.font.name = FONT_BODY
        run.font.color.rgb = _rgb(COLOR_BLACK)
        p.space_after = Pt(14)
